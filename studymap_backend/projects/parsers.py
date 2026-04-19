import os
import io
import fitz
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
import pytesseract


class DocumentParser:
    def detect_file_type(self, filename):
        ext = os.path.splitext(filename)[1].lower()
        if ext == '.pdf':
            return 'pdf'
        elif ext in ['.docx', '.doc']:
            return 'docx'
        elif ext == '.txt':
            return 'txt'
        elif ext in ['.pptx', '.ppt']:
            return 'pptx'
        raise ValueError(f'Unsupported file type: {ext}')

    def parse(self, file_path, file_type):
        if file_type == 'pdf':
            return self.parse_pdf(file_path)
        elif file_type == 'docx':
            return self.parse_docx(file_path)
        elif file_type == 'txt':
            return self.parse_txt(file_path)
        elif file_type == 'pptx':
            return self.parse_pptx(file_path)
        raise ValueError(f'Unsupported file type: {file_type}')

    def parse_pdf(self, file_path):
        text_parts = []
        page_count = 0

        try:
            doc = fitz.open(file_path)
            page_count = len(doc)

            for page_num, page in enumerate(doc):
                blocks = page.get_text('blocks')
                if blocks:
                    page_text = '\n\n'.join(block[4] for block in blocks if block[4].strip())
                    if page_text.strip():
                        text_parts.append(page_text)
                else:
                    pix = page.get_pixmap(dpi=300)
                    img_data = pix.tobytes('png')
                    img = Image.open(io.BytesIO(img_data))
                    ocr_text = pytesseract.image_to_string(img)
                    if ocr_text.strip():
                        text_parts.append(ocr_text)

                if page_num < len(doc) - 1:
                    text_parts.append('')

            doc.close()
        except Exception as e:
            if 'tesseract' in str(e).lower():
                raise RuntimeError('OCR (tesseract) is not installed. Install with: brew install tesseract')
            raise RuntimeError(f'PDF parsing error: {str(e)}')

        full_text = '\n\n'.join(text_parts)
        return full_text, page_count

    def parse_docx(self, file_path):
        try:
            doc = DocxDocument(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            full_text = '\n\n'.join(paragraphs)
            return full_text
        except Exception as e:
            raise RuntimeError(f'DOCX parsing error: {str(e)}')

    def parse_txt(self, file_path):
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            raise RuntimeError(f'Text file parsing error: {str(e)}')

    def parse_pptx(self, file_path):
        try:
            prs = Presentation(file_path)
            slides_count = len(prs.slides)

            text_parts = []
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text_frame.text.strip():
                                    slide_text.append(cell.text_frame.text)
                if slide_text:
                    text_parts.append('\n'.join(slide_text))

            full_text = '\n\n---\n\n'.join(text_parts)
            return full_text, slides_count
        except Exception as e:
            raise RuntimeError(f'PPTX parsing error: {str(e)}')